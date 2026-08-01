import pytest

from deccan_convert.ir import DocumentIR, Metadata
from deccan_convert.readers.md_reader import read_md
from deccan_convert.writers.html_writer import render_html, write_html
from deccan_convert.writers.md_writer import render_md
from deccan_convert.writers.docx_writer import write_docx

BANNED_FONTS = (
    "Helvetica", "Arial", "Calibri", "Verdana", "Times New Roman",
    "Garamond", "Georgia", "Courier New", "Lucida Console",
)


def _ir(sample_md):
    return read_md(sample_md)


class TestHtmlWriter:
    def test_all_slots_filled(self, sample_md):
        html = render_html(_ir(sample_md))
        assert "{{" not in html and "}}" not in html

    def test_print_contract_css_present(self, sample_md):
        html = render_html(_ir(sample_md))
        assert '"Deccan Fine Chemicals · Confidential"' in html
        assert "@page :first" in html
        assert "@page end-of-doc" in html
        assert "size: Letter" in html

    def test_light_only_rendering_contract(self, sample_md):
        """Previews that paint a dark canvas must not darken the document.

        The Claude iOS in-app viewer (and comparable webviews) render an
        embedded page against dark chrome; with the canvas left to a bare
        `body { background }` the dark stone text lands on a dark background.
        """
        html = render_html(_ir(sample_md))
        assert '<meta name="color-scheme" content="light">' in html
        assert "color-scheme: light only;" in html
        # Canvas pinned above the specificity of an injected html/body rule.
        assert ":root { background-color: var(--stone-50) !important; }" in html
        assert "html body {\n      background-color: var(--stone-50) !important;" in html
        assert "@media (prefers-color-scheme: dark) {" in html
        # ...and print still gets pure white paper over those rules.
        assert ":root { background-color: var(--paper) !important; }" in html
        assert "html body {\n        background-color: var(--paper) !important;" in html

    def test_metadata_escaped(self, sample_md):
        ir = _ir(sample_md)
        ir.metadata.title = 'A <b>"bold"</b> & risky title'
        html = render_html(ir)
        assert "<b>" not in html.split("<body>")[1].split("<main")[0]
        assert "&amp; risky" in html

    def test_required_fields_enforced(self):
        ir = DocumentIR(metadata=Metadata(title="X"), body_html="<p>b</p>")
        with pytest.raises(ValueError, match="document type"):
            render_html(ir)


class TestMdWriter:
    def test_roundtrip(self, sample_md, tmp_path):
        md_text = render_md(_ir(sample_md))
        assert md_text.startswith("---\n")
        assert "title: Solvent Recovery Audit" in md_text
        assert "# Findings" in md_text
        assert "| Unit |" in md_text
        # And it reads back cleanly.
        path = tmp_path / "again.md"
        path.write_text(md_text, encoding="utf-8")
        ir2 = read_md(path)
        assert ir2.metadata.title == "Solvent Recovery Audit"
        assert ir2.body_html.count('<section class="section">') == 2


class TestDocxWriter:
    @pytest.fixture
    def built(self, sample_md, tmp_path):
        import docx

        path = write_docx(_ir(sample_md), tmp_path / "out.docx")
        return docx.Document(str(path))

    def test_three_section_structure_preserved(self, built):
        assert len(built.sections) == 3

    def test_cover_filled(self, built):
        texts = [p.text for p in built.paragraphs]
        assert "Solvent Recovery Audit" in texts
        assert "REPORT" in texts
        meta_table = built.tables[0]
        values = [c.text for c in meta_table.rows[1].cells]
        assert values == ["Report", "QHSE Team", values[2], "0.3", "Internal"]

    def test_body_content_styled(self, built):
        styles = {p.text: p.style.name for p in built.paragraphs if p.text.strip()}
        assert styles.get("01  Findings") == "Heading 1"
        assert styles.get("Unit-level results") == "Heading 2"
        assert (
            styles.get("The recovery loop meets the target in four of five units.")
            == "Lead"
        )
        assert (
            styles.get("Unit B-2 requires a condenser retrofit before Q3.")
            == "Callout Default"
        )
        assert styles.get("deccan-convert audit.md -o audit.pdf") == "Code Block"

    def test_footer_contract_intact(self, built):
        # Body section keeps the template footer; cover/end sections do not
        # inherit it (different-first-page + unlinked footers).
        body_footer = built.sections[1].footer
        assert not body_footer.is_linked_to_previous
        footer_xml = body_footer.paragraphs[0]._p.xml
        assert "Cascadia Mono" in footer_xml

    def test_no_banned_fonts(self, sample_md, tmp_path):
        import zipfile

        path = write_docx(_ir(sample_md), tmp_path / "fonts.docx")
        document_xml = zipfile.ZipFile(path).read("word/document.xml").decode()
        for face in BANNED_FONTS:
            assert f'w:ascii="{face}"' not in document_xml

    def test_sample_content_removed(self, built):
        texts = " ".join(p.text for p in built.paragraphs)
        assert "Replace this content" not in texts
        assert "Document title" not in texts


class TestXlsxWriter:
    def test_data_preserved_and_styled(self, sample_xlsx, tmp_path):
        import openpyxl

        from deccan_convert.writers.xlsx_writer import restyle_xlsx

        out, warnings = restyle_xlsx(sample_xlsx, tmp_path / "out.xlsx")
        wb = openpyxl.load_workbook(str(out))
        ws = wb["Costs"]
        assert ws["A2"].value == "Reactor seal"
        assert ws["D2"].value == "=B2*C2"  # formulas verbatim
        assert ws["B2"].number_format == "#,##0"
        assert "A5:D5" in [str(r) for r in ws.merged_cells.ranges]
        # Header recipe from the .xltx: mono, stone-500, stone-100 fill.
        assert ws["A1"].font.name == "Cascadia Mono"
        assert ws["A1"].font.bold
        assert ws["A1"].fill.fgColor.rgb.endswith("F5F5F4")
        assert ws["A2"].font.name == "Segoe UI Variable Text"
        assert ws.freeze_panes == "A2"

    def test_input_never_mutated(self, sample_xlsx, tmp_path):
        import openpyxl

        from deccan_convert.writers.xlsx_writer import restyle_xlsx

        before = sample_xlsx.read_bytes()
        restyle_xlsx(sample_xlsx, tmp_path / "out.xlsx")
        assert sample_xlsx.read_bytes() == before


class TestPptxWriter:
    def test_restyle(self, sample_pptx, tmp_path):
        from pptx import Presentation

        from deccan_convert.writers.pptx_writer import restyle_pptx

        out, warnings = restyle_pptx(sample_pptx, tmp_path / "out.pptx")
        deck = Presentation(str(out))
        slides = list(deck.slides)
        # 3 source slides + generated end slide.
        assert len(slides) == 4
        cover_text = " ".join(
            sh.text_frame.text for sh in slides[0].shapes if sh.has_text_frame
        )
        assert "Plant Expansion" in cover_text
        assert "Deccan Fine Chemicals" in cover_text
        section_text = " ".join(
            sh.text_frame.text for sh in slides[1].shapes if sh.has_text_frame
        )
        assert "SECTION 01" in section_text and "Timeline" in section_text
        assert slides[2].notes_slide.notes_text_frame.text == "Mention the vendor delay."
        end_text = " ".join(
            sh.text_frame.text for sh in slides[3].shapes if sh.has_text_frame
        )
        assert "deccanchemicals.com" in end_text

    def test_fonts_are_tokens(self, sample_pptx, tmp_path):
        from pptx import Presentation

        from deccan_convert.writers.pptx_writer import restyle_pptx

        out, _ = restyle_pptx(sample_pptx, tmp_path / "out.pptx")
        deck = Presentation(str(out))
        seen = set()
        for slide in deck.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            seen.add(run.font.name)
        assert seen <= {
            "Segoe UI Variable Display",
            "Segoe UI Variable Text",
            "Cascadia Mono",
        }
        for face in BANNED_FONTS:
            assert face not in seen
