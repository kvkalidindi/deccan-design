"""Tests for the template selector, export kit, and logo option."""

import subprocess
import sys
from pathlib import Path

import pytest

from deccan_convert.convert import convert
from deccan_convert.ir import Metadata
from deccan_convert.kit import export_kit
from deccan_convert.readers.md_reader import read_md
from deccan_convert.writers.docx_writer import WORD_TEMPLATES, write_docx
from deccan_convert.writers.html_writer import render_html


class TestTemplateSelector:
    @pytest.mark.parametrize("template", sorted(WORD_TEMPLATES))
    def test_every_flavor_builds_with_positional_cover(self, template, sample_md, tmp_path):
        import docx

        ir = read_md(sample_md)
        path = write_docx(ir, tmp_path / f"{template}.docx", template=template)
        d = docx.Document(str(path))
        assert len(d.sections) == 3
        texts = [p.text for p in d.paragraphs if p.text.strip()]
        # Positional fill: eyebrow, title, subtitle right after the wordmark.
        assert texts[0] == "— Deccan Fine Chemicals"
        assert texts[1] == "REPORT"
        assert texts[2] == "Solvent Recovery Audit"
        # End-page classification stamped regardless of flavor placeholder.
        assert "INTERNAL · INTERNAL USE" in texts
        # Footer contract survives in every base.
        footer_xml = d.sections[1].footer.paragraphs[0]._p.xml
        assert "Cascadia Mono" in footer_xml

    def test_implied_document_type(self, tmp_path):
        md = tmp_path / "untyped.md"
        md.write_text(
            "---\ntitle: Reactor Spec\nauthor: Engineering\n---\n\n# Scope\n\nBody.\n",
            encoding="utf-8",
        )
        import docx

        ir = read_md(md)
        assert not ir.metadata.document_type
        path = write_docx(ir, tmp_path / "spec.docx", template="technical-spec")
        d = docx.Document(str(path))
        texts = [p.text for p in d.paragraphs if p.text.strip()]
        assert texts[1] == "SPECIFICATION"

    def test_unknown_template_rejected(self, sample_md, tmp_path):
        ir = read_md(sample_md)
        with pytest.raises(ValueError, match="Unknown Word template"):
            write_docx(ir, tmp_path / "x.docx", template="brochure")

    def test_template_rejected_for_non_docx_output(self, sample_md, tmp_path):
        with pytest.raises(ValueError, match="Word .* output only"):
            convert(
                sample_md, tmp_path / "out.html",
                metadata=Metadata(title="T", document_type="Report", prepared_by="X"),
                template="policy",
            )


class TestExportKit:
    def test_kit_contents(self, tmp_path):
        target = export_kit(tmp_path)
        assert target == tmp_path / "deccan-design-kit"
        expected = [
            "README.md",
            "templates/word/deccan-document.dotx",
            "templates/word/deccan-policy.dotx",
            "templates/excel/deccan-financial-model.xltx",
            "templates/powerpoint/deccan-customer-pitch.potx",
            "templates/gworkspace/deccan-gmail-signature.html",
            "templates/outlook/deccan-signature.htm",
            "skill/SKILL.md",
            "skill/references/tokens.md",
            "skill/assets/logo.b64.txt",
            "skill/assets/templates/document.html",
        ]
        for rel in expected:
            assert (target / rel).is_file(), f"missing from kit: {rel}"
        readme = (target / "README.md").read_text(encoding="utf-8")
        assert ".claude/skills" in readme
        assert "text wordmark" in readme

    def test_kit_skill_carries_the_fetch_first_rule(self, tmp_path):
        """The kit equips offline endpoints — it must still teach fetch-first.

        An endpoint equipped from `--export-kit` gets a skill copy that is
        frozen at the binary's build date. If that copy tells a session the
        bundled template is good enough, the endpoint produces dark-on-dark
        documents forever. The rule that sends the session to the canonical
        URL is what bounds the staleness, so it ships in the kit or the kit
        is not fit to equip anything.
        """
        skill = (export_kit(tmp_path) / "skill" / "SKILL.md").read_text(encoding="utf-8")
        assert "## Fetching the template — hard rule" in skill
        assert "Fetch it at build time, every time" in skill
        assert "## The rendering invariant" in skill
        for marker in ("color-scheme: light only;",
                       ":root { background-color: var(--stone-50) !important; }"):
            assert marker in skill, f"invariant does not list {marker}"

    def test_kit_template_is_render_safe(self, tmp_path):
        """The template inside the kit carries the contract it advertises."""
        tpl = (export_kit(tmp_path) / "skill" / "assets" / "templates"
               / "document.html").read_text(encoding="utf-8")
        for marker in ('<meta name="color-scheme" content="light">',
                       '<meta name="generator" content="deccan-design v2.0 · slot template ',
                       "color-scheme: light only;",
                       ":root { background-color: var(--stone-50) !important; }",
                       "@media (prefers-color-scheme: dark) {"):
            assert marker in tpl, f"kit template missing: {marker}"

    def test_kit_refuses_to_overwrite(self, tmp_path):
        export_kit(tmp_path)
        with pytest.raises(FileExistsError):
            export_kit(tmp_path)

    def test_cli_export_kit(self, tmp_path):
        proc = subprocess.run(
            [sys.executable, "-m", "deccan_convert", "--export-kit", str(tmp_path / "k")],
            capture_output=True, text=True, cwd=Path(__file__).parents[1],
        )
        assert proc.returncode == 0, proc.stderr
        assert (tmp_path / "k" / "deccan-design-kit" / "skill" / "SKILL.md").is_file()


class TestLogoOption:
    def test_html_logo_cover(self, sample_md):
        ir = read_md(sample_md)
        html = render_html(ir, logo=True)
        assert "data:image/png;base64," in html
        assert 'class="mark-logo"' in html
        # The cover's text mark element is replaced (the CSS rule remains);
        # the end page keeps its text mark.
        assert '<div class="mark-text">' not in html
        assert 'class="end-mark"' in html

    def test_html_default_stays_text(self, sample_md):
        ir = read_md(sample_md)
        html = render_html(ir)
        assert "data:image/png;base64," not in html
        assert 'class="mark-text"' in html

    def test_docx_logo_embeds_image(self, sample_md, tmp_path):
        import zipfile

        ir = read_md(sample_md)
        path = write_docx(ir, tmp_path / "logo.docx", logo=True)
        names = zipfile.ZipFile(path).namelist()
        assert any(n.startswith("word/media/") for n in names)

    def test_pptx_logo_on_cover(self, sample_pptx, tmp_path):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        from deccan_convert.writers.pptx_writer import restyle_pptx

        out, _ = restyle_pptx(sample_pptx, tmp_path / "out.pptx", logo=True)
        deck = Presentation(str(out))
        cover = list(deck.slides)[0]
        assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in cover.shapes)
