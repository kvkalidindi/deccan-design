from __future__ import annotations

import os
from pathlib import Path

import pytest

# Browser-marked tests render through headless Chromium; on root/container
# test hosts the Chromium sandbox is unavailable, so opt into the gated
# --no-sandbox retry for the test session only.
os.environ.setdefault("DECCAN_CONVERT_ALLOW_NO_SANDBOX", "1")

FIXTURES = Path(__file__).parent / "fixtures"

SAMPLE_MD = """\
---
title: Solvent Recovery Audit
subtitle: Q2 findings for the Hyderabad plant
type: Report
author: QHSE Team
version: "0.3"
classification: Internal
---

# Findings

The recovery loop meets the target in four of five units.

## Unit-level results

| Unit | Recovery | Target |
|------|----------|--------|
| A-1  | 94%      | 92%    |
| B-2  | 89%      | 92%    |

> Unit B-2 requires a condenser retrofit before Q3.

# Recommendations

Replace the B-2 condenser. See [the plan](https://example.com/plan).

```bash
deccan-convert audit.md -o audit.pdf
```
"""


@pytest.fixture
def sample_md(tmp_path: Path) -> Path:
    path = tmp_path / "sample.md"
    path.write_text(SAMPLE_MD, encoding="utf-8")
    return path


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    import docx

    d = docx.Document()
    d.core_properties.title = "Vendor Assessment"
    d.core_properties.author = "Procurement"
    d.add_heading("Scope", level=1)
    d.add_paragraph("The assessment covers three vendors.")
    d.add_heading("Criteria", level=2)
    p = d.add_paragraph("Cost, ")
    p.add_run("quality").bold = True
    p.add_run(", and lead time.")
    d.add_paragraph("Vendor one", style="List Bullet")
    d.add_paragraph("Vendor two", style="List Bullet")
    table = d.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Vendor"
    table.rows[0].cells[1].text = "Score"
    table.rows[1].cells[0].text = "Acme"
    table.rows[1].cells[1].text = "87"
    path = tmp_path / "sample.docx"
    d.save(str(path))
    return path


@pytest.fixture
def sample_html(tmp_path: Path) -> Path:
    path = tmp_path / "sample.html"
    path.write_text(
        """<!DOCTYPE html>
<html><head><title>Maintenance Notice</title>
<meta name="author" content="Plant Ops"></head>
<body>
<h1>Maintenance Notice</h1>
<h2>Window</h2>
<p>The line stops <b>Friday</b> at 18:00.</p>
<script>alert('x')</script>
<ul><li>Drain loop</li><li>Lock out panel</li></ul>
</body></html>
""",
        encoding="utf-8",
    )
    return path


@pytest.fixture
def sample_xlsx(tmp_path: Path) -> Path:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Costs"
    ws.append(["Item", "Qty", "Unit", "Total"])
    ws.append(["Reactor seal", 4, 120, "=B2*C2"])
    ws.append(["Filter mesh", 10, 15, "=B3*C3"])
    ws["A5"] = "Note"
    ws.merge_cells("A5:D5")
    ws["B2"].number_format = "#,##0"
    wb.create_sheet("Empty")
    path = tmp_path / "sample.xlsx"
    wb.save(str(path))
    return path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    from pptx import Presentation

    p = Presentation()
    s = p.slides.add_slide(p.slide_layouts[0])
    s.shapes.title.text = "Plant Expansion"
    s.placeholders[1].text = "Board update, July 2026"
    s2 = p.slides.add_slide(p.slide_layouts[2])
    s2.shapes.title.text = "Timeline"
    s3 = p.slides.add_slide(p.slide_layouts[1])
    s3.shapes.title.text = "Milestones"
    body = s3.placeholders[1].text_frame
    body.text = "Phase one complete"
    q = body.add_paragraph()
    q.text = "Permits filed"
    q.level = 1
    s3.notes_slide.notes_text_frame.text = "Mention the vendor delay."
    path = tmp_path / "sample.pptx"
    p.save(str(path))
    return path


@pytest.fixture
def sample_pdf() -> Path:
    return FIXTURES / "sample.pdf"


def browser_available() -> bool:
    from deccan_convert.writers.pdf_writer import BrowserNotFound, find_browser

    try:
        find_browser()
        return True
    except BrowserNotFound:
        return False
