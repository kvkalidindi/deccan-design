"""PPTX restyler: source .pptx -> restyled deck matching deccan-deck.pptx.

The Deccan deck's design is carried as direct formatting (the bundled deck
targets OS-native fonts, not a theme), so this restyler rebuilds each source
slide inside a copy of the bundled deck using the same geometry and token
typography as its sample slides: cover, section breaks, content slides,
tables, and a closing end slide. Slide text, bullet structure, tables,
pictures, and speaker notes carry over; SmartArt and embedded charts cannot
be ported and are replaced with a placeholder note. Decks exported from
Google Slides (File > Download > .pptx) take the same path.
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Callable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from deccan_convert.assets import asset_path
from deccan_convert.ir import current_month_year

# Tokens (skill/references/tokens.md).
DECCAN_BLUE = RGBColor(0x16, 0x49, 0x99)
STONE_100 = RGBColor(0xF5, 0xF5, 0xF4)
STONE_200 = RGBColor(0xE7, 0xE5, 0xE4)
STONE_500 = RGBColor(0x78, 0x71, 0x6C)
STONE_700 = RGBColor(0x44, 0x40, 0x3C)
STONE_800 = RGBColor(0x29, 0x25, 0x24)
STONE_900 = RGBColor(0x1C, 0x19, 0x17)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)

SANS_DISPLAY = "Segoe UI Variable Display"
SANS_TEXT = "Segoe UI Variable Text"
MONO = "Cascadia Mono"

_CHART_TYPES = (MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT)


def restyle_pptx(
    source: Path,
    path: Path,
    log: Callable[[str], None] | None = None,
    logo: bool = False,
) -> tuple[Path, list[str]]:
    say = log or (lambda _msg: None)
    warnings: list[str] = []

    src = Presentation(str(source))
    deck = Presentation(str(asset_path("deccan-deck.pptx")))
    blank_layout = _blank_layout(deck)
    _delete_all_slides(deck)

    slides = list(src.slides)
    say(f"Restyling {len(slides)} slide(s)")

    for index, slide in enumerate(slides):
        content = _harvest(slide, index, warnings)
        new_slide = deck.slides.add_slide(blank_layout)
        if index == 0 and content["looks_like_cover"]:
            _build_cover(deck, new_slide, content, logo=logo)
        elif content["is_section_break"]:
            _build_section(deck, new_slide, content)
        else:
            _build_content(deck, new_slide, content, warnings)
        if content["notes"]:
            new_slide.notes_slide.notes_text_frame.text = content["notes"]

    _build_end_slide(deck, deck.slides.add_slide(blank_layout))

    deck.save(str(path))
    return path, warnings


# --- source harvesting -------------------------------------------------------


def _harvest(slide, index: int, warnings: list[str]) -> dict:
    title = ""
    bodies: list[list[tuple[int, str]]] = []  # text frames: [(indent level, text)]
    tables: list[list[list[str]]] = []
    pictures: list[bytes] = []
    dropped: list[str] = []

    title_shape = slide.shapes.title
    if title_shape is not None and title_shape.has_text_frame:
        title = title_shape.text_frame.text.strip()

    for shape in slide.shapes:
        if title_shape is not None and shape.shape_id == title_shape.shape_id:
            continue
        if shape.has_table:
            table = shape.table
            tables.append(
                [
                    [cell.text.strip() for cell in row.cells]
                    for row in table.rows
                ]
            )
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                pictures.append(shape.image.blob)
            except Exception:
                dropped.append("picture")
        elif shape.shape_type in _CHART_TYPES or shape.has_chart:
            dropped.append("chart/embedded object")
        elif shape.has_text_frame:
            paragraphs = [
                (para.level, para.text.strip())
                for para in shape.text_frame.paragraphs
                if para.text.strip()
            ]
            if paragraphs:
                if not title and shape.is_placeholder and shape.placeholder_format.type in (
                    PP_PLACEHOLDER.TITLE,
                    PP_PLACEHOLDER.CENTER_TITLE,
                ):
                    title = paragraphs[0][1]
                else:
                    bodies.append(paragraphs)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            dropped.append("grouped shapes")

    if dropped:
        warnings.append(
            f"pptx: slide {index + 1}: {', '.join(sorted(set(dropped)))} could not "
            "be ported and was replaced with a placeholder note."
        )

    body_text_len = sum(len(t) for body in bodies for _, t in body)
    subtitle = ""
    if bodies and index == 0:
        subtitle = " ".join(t for _, t in bodies[0])[:120]

    return {
        "title": title,
        "bodies": bodies,
        "tables": tables,
        "pictures": pictures,
        "dropped": dropped,
        "notes": (
            slide.notes_slide.notes_text_frame.text.strip()
            if slide.has_notes_slide
            else ""
        ),
        "looks_like_cover": index == 0 and body_text_len < 300 and not tables,
        "is_section_break": (
            index > 0 and bool(title) and not bodies and not tables and not pictures
        ),
        "subtitle": subtitle,
    }


# --- deck plumbing -----------------------------------------------------------


def _blank_layout(deck):
    for layout in deck.slide_masters[0].slide_layouts:
        if layout.name == "Blank":
            return layout
    return deck.slide_masters[0].slide_layouts[-1]


def _delete_all_slides(deck) -> None:
    xml_slides = deck.slides._sldIdLst
    for sld_id in list(xml_slides):
        rel_id = sld_id.rId
        deck.part.drop_rel(rel_id)
        xml_slides.remove(sld_id)


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.word_wrap = True
    return box


def _set_run(run, text, font=SANS_TEXT, size=18, bold=False, color=STONE_800):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _blue_rule(slide, left, top, width=0.6):
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = DECCAN_BLUE
    rule.line.fill.background()
    return rule


# --- slide builders (geometry mirrors the bundled deck's sample slides) -------


def _build_cover(deck, slide, content, logo: bool = False) -> None:
    if logo:
        # Graphical wordmark from the bundled asset (never fetched); the
        # text mark is the sanctioned fallback everywhere else.
        from pptx.util import Inches

        slide.shapes.add_picture(
            str(asset_path("logo.png")), Inches(0.6), Inches(2.7), height=Inches(0.5)
        )
    else:
        box = _textbox(slide, 0.6, 2.8, 11.0, 0.4)
        _set_run(box.text_frame.paragraphs[0].add_run(), "Deccan Fine Chemicals",
                 SANS_TEXT, 12, True, STONE_900)
    _blue_rule(slide, 0.6, 3.5)
    box = _textbox(slide, 0.6, 3.7, 11.0, 1.4)
    _set_run(box.text_frame.paragraphs[0].add_run(),
             content["title"] or "Presentation", SANS_DISPLAY, 44, False, DECCAN_BLUE)
    if content["subtitle"]:
        box = _textbox(slide, 0.6, 5.0, 11.0, 0.6)
        _set_run(box.text_frame.paragraphs[0].add_run(), content["subtitle"],
                 SANS_TEXT, 18, False, STONE_700)
    box = _textbox(slide, 0.6, 6.4, 8.0, 0.4)
    meta = f"{current_month_year().upper()} · VERSION 1.0 · CONFIDENTIAL"
    _set_run(box.text_frame.paragraphs[0].add_run(), meta, MONO, 10, True, STONE_500)


_SECTION_COUNTER_KEY = "_deccan_section_no"


def _build_section(deck, slide, content) -> None:
    number = getattr(deck, _SECTION_COUNTER_KEY, 0) + 1
    setattr(deck, _SECTION_COUNTER_KEY, number)
    box = _textbox(slide, 0.6, 0.4, 4.0, 0.4)
    _set_run(box.text_frame.paragraphs[0].add_run(), f"SECTION {number:02d}",
             MONO, 11, True, STONE_500)
    box = _textbox(slide, 0.6, 3.0, 12.0, 1.4)
    _set_run(box.text_frame.paragraphs[0].add_run(), content["title"],
             SANS_DISPLAY, 44, False, DECCAN_BLUE)
    _blue_rule(slide, 0.6, 4.6)


def _build_content(deck, slide, content, warnings: list[str]) -> None:
    top = 1.4
    if content["title"]:
        box = _textbox(slide, 0.6, 0.5, 12.0, 0.8)
        _set_run(box.text_frame.paragraphs[0].add_run(), content["title"],
                 SANS_TEXT, 24, True, STONE_900)

    bodies = content["bodies"]
    if bodies:
        columns = min(len(bodies), 3) if len(bodies) > 1 else 1
        if columns == 1:
            merged = [p for body in bodies for p in body]
            _fill_body(_textbox(slide, 0.6, top, 12.0, 5.0), merged, size=18)
        else:
            width = {2: 5.8, 3: 3.8}[columns]
            step = {2: 6.4, 3: 4.1}[columns]
            size = {2: 16, 3: 14}[columns]
            for i, body in enumerate(bodies[:columns]):
                _fill_body(
                    _textbox(slide, 0.6 + i * step, top, width, 5.0), body, size=size
                )
            extra = [p for body in bodies[columns:] for p in body]
            if extra:
                warnings.append(
                    "pptx: more than three text blocks on one slide; the "
                    "overflow was merged into the last column."
                )
                _fill_body(
                    _textbox(slide, 0.6 + (columns - 1) * step, top, width, 5.0),
                    bodies[columns - 1] + extra,
                    size=size,
                )
        top = 6.4

    for table_data in content["tables"]:
        _build_table(slide, table_data)
    if content["dropped"]:
        box = _textbox(slide, 0.6, 6.6, 12.0, 0.5)
        _set_run(
            box.text_frame.paragraphs[0].add_run(),
            "[Chart or embedded object omitted — recreate from the source deck]",
            MONO, 10, False, STONE_500,
        )
    _place_pictures(slide, content["pictures"])


def _fill_body(box, paragraphs: list[tuple[int, str]], size: int) -> None:
    tf = box.text_frame
    for i, (level, text) in enumerate(paragraphs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = min(level, 4)
        run = para.add_run()
        bullet = "· " if level > 0 else ""
        _set_run(run, bullet + text, SANS_TEXT, max(size - 2 * level, 10),
                 False, STONE_800)
        para.space_after = Pt(6)


def _build_table(slide, data: list[list[str]]) -> None:
    rows, cols = len(data), max(len(r) for r in data)
    if not rows or not cols:
        return
    shape = slide.shapes.add_table(
        rows, cols, Inches(0.6), Inches(1.6), Inches(12.0), Inches(min(0.4 * rows, 4.8))
    )
    table = shape.table
    table.first_row = False
    table.horz_banding = False
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            text = data[r][c] if c < len(data[r]) else ""
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            if r == 0:
                _set_run(run, text.upper(), MONO, 10, True, STONE_700)
                cell.fill.solid()
                cell.fill.fore_color.rgb = STONE_100
            else:
                _set_run(run, text, SANS_TEXT, 12, False, STONE_800)
                cell.fill.solid()
                cell.fill.fore_color.rgb = PAPER
            para.alignment = PP_ALIGN.LEFT


def _place_pictures(slide, pictures: list[bytes]) -> None:
    left = 0.6
    for blob in pictures[:3]:
        try:
            picture = slide.shapes.add_picture(
                io.BytesIO(blob), Inches(left), Inches(1.6), height=Inches(4.5)
            )
            # Keep within the slide's live area.
            max_right = Inches(12.7)
            if picture.left + picture.width > max_right:
                picture.width = Emu(int(max_right - picture.left))
            left += 4.3
        except Exception:
            continue


def _build_end_slide(deck, slide) -> None:
    width_in = deck.slide_width / 914400
    box = _textbox(slide, 0.0, 3.2, width_in, 0.5)
    para = box.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    _set_run(para.add_run(), "Deccan Fine Chemicals", SANS_TEXT, 21, True, STONE_900)
    _blue_rule(slide, width_in / 2 - 0.3, 3.9)
    box = _textbox(slide, 0.0, 4.2, width_in, 0.4)
    para = box.text_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    _set_run(para.add_run(), "deccanchemicals.com · Hyderabad, India",
             SANS_TEXT, 12, False, STONE_700)
