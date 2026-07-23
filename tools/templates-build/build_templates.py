#!/usr/bin/env python3
"""Remediate the deccan-design template suite for strict token adherence.

Patches the Office template sources in place (theme, styles, footers, sheet
views, master placeholders, docProps), regenerates the two specialised
PowerPoint templates from the patched deck, and re-derives the Google
Workspace variants from the patched Office sources.

Deterministic: fixed zip timestamps (2020-01-01), original entry order,
idempotent patches — re-running on already-patched sources is a byte-stable
no-op. Never hand-edit the derived files (gworkspace trio, the two
specialised .potx); rerun this script instead.

Usage:
    python build_templates.py           # patch + derive
    python build_templates.py --check   # exit 1 if any output would change
"""

from __future__ import annotations

import argparse
import io
import re
import sys
import zipfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import deccan_tokens as T

HERE = Path(__file__).resolve().parent
REPO = HERE.parent.parent
TEMPLATES = REPO / "templates"

CHECK_MODE = False
_changed: list[str] = []


# --- deterministic zip plumbing ---------------------------------------------


def read_parts(path_or_bytes) -> tuple[list[str], dict[str, bytes]]:
    src = io.BytesIO(path_or_bytes) if isinstance(path_or_bytes, bytes) else str(path_or_bytes)
    with zipfile.ZipFile(src) as zf:
        order = [i.filename for i in zf.infolist()]
        parts = {name: zf.read(name) for name in order}
    return order, parts


def zip_bytes(order: list[str], parts: dict[str, bytes]) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for name in order:
            zi = zipfile.ZipInfo(name, date_time=(2020, 1, 1, 0, 0, 0))
            zi.compress_type = zipfile.ZIP_DEFLATED
            zf.writestr(zi, parts[name])
    return buf.getvalue()


def emit(path: Path, data: bytes) -> None:
    """Write (or in --check mode, compare) an output file."""
    rel = path.relative_to(REPO)
    if path.is_file() and path.read_bytes() == data:
        print(f"  ok       {rel}")
        return
    if CHECK_MODE:
        _changed.append(str(rel))
        print(f"  DRIFT    {rel}")
    else:
        path.write_bytes(data)
        print(f"  patched  {rel}")


def sub_once(pattern: str, replacement: str, text: str, label: str, flags=re.DOTALL) -> str:
    new, n = re.subn(pattern, replacement, text, count=1, flags=flags)
    if n != 1:
        raise SystemExit(f"patch anchor not found: {label}")
    return new


# --- shared part patches ------------------------------------------------------


def patch_theme(xml: str) -> str:
    xml = sub_once(r"<a:clrScheme .*?</a:clrScheme>", T.CLR_SCHEME, xml, "clrScheme")
    xml = sub_once(r"<a:fontScheme .*?</a:fontScheme>", T.FONT_SCHEME, xml, "fontScheme")
    return xml


def patch_core(xml: str, title: str) -> str:
    """Tolerates both core.xml shapes: the full python-docx/pptx layout and
    openpyxl's minimal one (inline namespaces, no title element)."""
    dc_ns = ' xmlns:dc="http://purl.org/dc/elements/1.1/"' if "xmlns:dc=" not in xml.split(">", 1)[0] else ""

    if re.search(r"<dc:title", xml):
        xml = re.sub(r"<dc:title([^>]*)/>|<dc:title([^>]*)>.*?</dc:title>",
                     lambda m: f"<dc:title{m.group(1) or m.group(2) or ''}>{title}</dc:title>",
                     xml, count=1, flags=re.DOTALL)
    else:
        xml = sub_once(r"(<cp:coreProperties[^>]*>)",
                       rf"\g<1><dc:title{dc_ns}>{title}</dc:title>", xml, "coreProperties open")

    xml = re.sub(r"<dc:creator([^>]*)/>|<dc:creator([^>]*)>.*?</dc:creator>",
                 lambda m: f"<dc:creator{m.group(1) or m.group(2) or ''}>{T.COMPANY}</dc:creator>",
                 xml, count=1, flags=re.DOTALL)

    if re.search(r"<cp:lastModifiedBy", xml):
        xml = re.sub(r"<cp:lastModifiedBy\s*/>|<cp:lastModifiedBy>.*?</cp:lastModifiedBy>",
                     f"<cp:lastModifiedBy>{T.COMPANY}</cp:lastModifiedBy>", xml, count=1, flags=re.DOTALL)
    else:
        xml = sub_once(r"</cp:coreProperties>",
                       f"<cp:lastModifiedBy>{T.COMPANY}</cp:lastModifiedBy></cp:coreProperties>",
                       xml, "coreProperties close")

    xml = re.sub(r"<dc:description\s*/>|<dc:description>.*?</dc:description>",
                 "<dc:description>deccan-design v2.0 template</dc:description>",
                 xml, count=1, flags=re.DOTALL)
    xml = re.sub(r'(<dcterms:created[^>]*xsi:type="dcterms:W3CDTF"[^>]*>).*?(</dcterms:created>)',
                 r"\g<1>2020-01-01T00:00:00Z\g<2>", xml, count=1, flags=re.DOTALL)
    xml = re.sub(r'(<dcterms:modified[^>]*xsi:type="dcterms:W3CDTF"[^>]*>).*?(</dcterms:modified>)',
                 r"\g<1>2020-01-01T00:00:00Z\g<2>", xml, count=1, flags=re.DOTALL)
    return xml


def patch_app(xml: str) -> str:
    if "<Company>" in xml:
        return re.sub(r"<Company>.*?</Company>", f"<Company>{T.COMPANY}</Company>", xml, count=1)
    return sub_once(r"</Properties>",
                    f"<Company>{T.COMPANY}</Company></Properties>", xml, "app </Properties>")


# --- Word --------------------------------------------------------------------

WORD_TITLES = {
    "deccan-document.dotx": "Deccan Document",
    "deccan-technical-spec.dotx": "Deccan Technical Specification",
    "deccan-policy.dotx": "Deccan Policy",
    "deccan-customer-letter.dotx": "Deccan Customer Letter",
}

_FOOTER_RPR = (
    '<w:rPr><w:rFonts w:ascii="Cascadia Mono" w:hAnsi="Cascadia Mono"'
    ' w:cs="Cascadia Mono" w:eastAsia="Cascadia Mono"/>'
    '<w:color w:val="78716C"/><w:sz w:val="17"/></w:rPr>'
)

# Furniture rule 6: left "Deccan Fine Chemicals · Confidential", right bare
# PAGE number, Cascadia Mono 8.5pt stone-500. The tab stop is flipped to a
# right tab at the margin edge (the template shipped it as a left tab).
FOOTER_PARAGRAPH = (
    "<w:p><w:pPr><w:pStyle w:val=\"Footer\"/>"
    "<w:tabs><w:tab w:pos=\"9936\" w:val=\"right\"/></w:tabs></w:pPr>"
    f"<w:r>{_FOOTER_RPR}<w:t xml:space=\"preserve\">{T.FOOTER_TEXT}</w:t></w:r>"
    f"<w:r>{_FOOTER_RPR}<w:tab/></w:r>"
    f"<w:r>{_FOOTER_RPR}<w:fldChar w:fldCharType=\"begin\"/></w:r>"
    f"<w:r>{_FOOTER_RPR}<w:instrText xml:space=\"preserve\"> PAGE </w:instrText></w:r>"
    f"<w:r>{_FOOTER_RPR}<w:fldChar w:fldCharType=\"separate\"/></w:r>"
    f"<w:r>{_FOOTER_RPR}<w:t>1</w:t></w:r>"
    f"<w:r>{_FOOTER_RPR}<w:fldChar w:fldCharType=\"end\"/></w:r>"
    "</w:p>"
)

_SEGOE_RFONTS = (
    '<w:rFonts w:ascii="Segoe UI Variable Text" w:hAnsi="Segoe UI Variable Text"'
    ' w:cs="Segoe UI Variable Text" w:eastAsia="Segoe UI Variable Text"/>'
)

# Heading 4 per the print scale: 11pt (sz 22), w600 (bold), stone-900.
_H4_RPR = f"<w:rPr>{_SEGOE_RFONTS}<w:b/><w:bCs/><w:color w:val=\"1C1917\"/><w:sz w:val=\"22\"/></w:rPr>"
# Headings 5-9 stay latent but deterministic: 10.5pt bold stone-700.
_H59_RPR = f"<w:rPr>{_SEGOE_RFONTS}<w:b/><w:bCs/><w:color w:val=\"44403C\"/><w:sz w:val=\"21\"/></w:rPr>"


def patch_word_styles(xml: str, strict: bool = True) -> str:
    # Heading 4-9: replace each style's rPr wholesale with the token shape.
    for style_id, rpr in [("Heading4", _H4_RPR)] + [(f"Heading{n}", _H59_RPR) for n in range(5, 10)]:
        block = re.search(rf'<w:style [^>]*w:styleId="{style_id}">.*?</w:style>', xml)
        if block is None:
            if strict:
                raise SystemExit(f"style not found: {style_id}")
            continue
        patched = re.sub(r"<w:rPr>.*?</w:rPr>", rpr, block.group(0), count=1)
        xml = xml.replace(block.group(0), patched, 1)

    # Banned Courier in the macro styles -> Cascadia Mono.
    xml = xml.replace('w:ascii="Courier"', 'w:ascii="Cascadia Mono"')
    xml = xml.replace('w:hAnsi="Courier"', 'w:hAnsi="Cascadia Mono"')

    # Gallery remap: every stale Office hex collapses to a token (single-accent
    # system: foreign accent hues all become Deccan Blue).
    def remap(match: re.Match) -> str:
        attr, hexval = match.group(1), match.group(2).upper()
        if hexval in T.ALLOWED_HEXES:
            return match.group(0)
        replacement = T.HEX_REMAP.get(hexval) or T.tint_fallback(hexval)
        return f'{attr}"{replacement}"'

    xml = re.sub(r'(w:(?:val|fill|color)=)"([0-9A-Fa-f]{6})"', remap, xml)

    # Post-patch assertion: only token hexes may remain.
    leftover = {
        h.upper()
        for h in re.findall(r'w:(?:val|fill|color)="([0-9A-Fa-f]{6})"', xml)
        if h.upper() not in T.ALLOWED_HEXES
    }
    if leftover:
        raise SystemExit(f"styles.xml still carries non-token hexes: {sorted(leftover)}")
    return xml


_BANNED_FONT_DECLS = ["Times New Roman", "Cambria", "Calibri", "Courier", "Arial"]


def patch_font_table(xml: str) -> str:
    """Drop declarations of banned faces. These entries are metadata, not
    usage, but the ban list allows no appearance at all; Word regenerates
    declarations on demand for fonts a document actually uses."""
    for face in _BANNED_FONT_DECLS:
        xml = re.sub(rf'<w:font w:name="{face}">.*?</w:font>', "", xml, flags=re.DOTALL)
    return xml


def patch_word(path: Path, title: str) -> None:
    order, parts = read_parts(path)
    parts["word/theme/theme1.xml"] = patch_theme(parts["word/theme/theme1.xml"].decode()).encode()
    parts["word/styles.xml"] = patch_word_styles(parts["word/styles.xml"].decode()).encode()
    # Word-2010 compat twin of styles.xml — same stale gallery (its docDefaults
    # even carry Courier); patch leniently, shapes may differ from styles.xml.
    if "word/stylesWithEffects.xml" in parts:
        parts["word/stylesWithEffects.xml"] = patch_word_styles(
            parts["word/stylesWithEffects.xml"].decode(), strict=False
        ).encode()
    parts["word/fontTable.xml"] = patch_font_table(parts["word/fontTable.xml"].decode()).encode()
    footer = parts["word/footer1.xml"].decode()
    footer = sub_once(r"<w:p>.*</w:p>", FOOTER_PARAGRAPH, footer, "footer1 paragraph")
    parts["word/footer1.xml"] = footer.encode()
    parts["docProps/core.xml"] = patch_core(parts["docProps/core.xml"].decode(), title).encode()
    parts["docProps/app.xml"] = patch_app(parts["docProps/app.xml"].decode()).encode()
    emit(path, zip_bytes(order, parts))


# --- Excel -------------------------------------------------------------------

EXCEL_TITLES = {
    "deccan-workbook.xltx": "Deccan Workbook",
    "deccan-comparison.xltx": "Deccan Comparison",
    "deccan-financial-model.xltx": "Deccan Financial Model",
}

# Excel footer codes: size codes are integer-only, so the 8.5pt token rounds
# to 9. &K = font colour (stone-500). Escaped for XML.
EXCEL_FOOTER = (
    "&amp;L&amp;\"Cascadia Mono,Regular\"&amp;9&amp;K78716C"
    "Deccan Fine Chemicals · Confidential"
    "&amp;R&amp;\"Cascadia Mono,Regular\"&amp;9&amp;K78716C&amp;P"
)


def patch_excel_sheet(xml: str) -> str:
    # Screen gridlines off.
    if "showGridLines" not in xml:
        xml = sub_once(r"<sheetView ", '<sheetView showGridLines="0" ', xml, "sheetView")

    # Freeze the header row when row 1 carries styled header cells.
    has_header = re.search(r'<row r="1"><c r="A1" s="[1-9]', xml) is not None
    if has_header and "<pane " not in xml:
        xml = sub_once(
            r'<selection activeCell="A1" sqref="A1"/>',
            '<pane ySplit="1" topLeftCell="A2" activePane="bottomLeft" state="frozen"/>'
            '<selection pane="bottomLeft" activeCell="A2" sqref="A2"/>',
            xml,
            "sheetView selection",
        )

    # Letter page setup + the running footer, inserted after pageMargins per
    # the worksheet schema order.
    if "<headerFooter>" not in xml:
        xml = sub_once(
            r"(<pageMargins [^>]*/>)",
            r"\g<1>"
            '<pageSetup paperSize="1" orientation="portrait"/>'
            f"<headerFooter><oddFooter>{EXCEL_FOOTER}</oddFooter></headerFooter>",
            xml,
            "pageMargins",
        )
    return xml


def patch_excel(path: Path, title: str) -> None:
    order, parts = read_parts(path)
    parts["xl/theme/theme1.xml"] = patch_theme(parts["xl/theme/theme1.xml"].decode()).encode()

    styles = parts["xl/styles.xml"].decode()
    if '<font><name val="Calibri"/>' in styles:
        # Font id 0 — the workbook default every unstyled cell inherits.
        styles = sub_once(
            r'(<fonts count="\d+"><font><name val=")Calibri(")',
            r"\g<1>Segoe UI Variable Text\g<2>",
            styles,
            "default font Calibri",
        )
    parts["xl/styles.xml"] = styles.encode()

    for name in list(parts):
        if re.fullmatch(r"xl/worksheets/sheet\d+\.xml", name):
            parts[name] = patch_excel_sheet(parts[name].decode()).encode()

    parts["docProps/core.xml"] = patch_core(parts["docProps/core.xml"].decode(), title).encode()
    parts["docProps/app.xml"] = patch_app(parts["docProps/app.xml"].decode()).encode()
    emit(path, zip_bytes(order, parts))


# --- PowerPoint --------------------------------------------------------------

_PPT_MONO_DEFRPR = (
    '<a:defRPr sz="850"><a:solidFill><a:srgbClr val="78716C"/></a:solidFill>'
    '<a:latin typeface="Cascadia Mono"/></a:defRPr>'
)
_PPT_MONO_RPR = (
    '<a:rPr lang="en-US" sz="850"><a:solidFill><a:srgbClr val="78716C"/></a:solidFill>'
    '<a:latin typeface="Cascadia Mono"/></a:rPr>'
)


def _patch_master_placeholder(xml: str, ph_type: str, default_text: str | None = None) -> str:
    """Restyle a master footer-strip placeholder to mono 8.5pt stone-500."""
    block = re.search(rf'<p:sp>(?:(?!</p:sp>).)*?type="{ph_type}".*?</p:sp>', xml, re.DOTALL)
    if block is None:
        raise SystemExit(f"master placeholder not found: {ph_type}")
    patched = re.sub(r"<a:defRPr .*?</a:defRPr>", _PPT_MONO_DEFRPR, block.group(0), count=1)
    if default_text is not None:
        patched = patched.replace(
            '<a:p><a:endParaRPr lang="en-US"/></a:p>',
            f"<a:p><a:r>{_PPT_MONO_RPR}<a:t>{default_text}</a:t></a:r></a:p>",
            1,
        )
    return xml.replace(block.group(0), patched, 1)


def patch_ppt_master(xml: str) -> str:
    # Banned bullet font.
    xml = xml.replace('<a:buFont typeface="Arial"/>', '<a:buFont typeface="Segoe UI Variable Text"/>')
    xml = xml.replace('<a:latin typeface="Calibri"/>', '<a:latin typeface="+mn-lt"/>')
    xml = _patch_master_placeholder(xml, "ftr", default_text=T.FOOTER_TEXT)
    xml = _patch_master_placeholder(xml, "sldNum")
    xml = _patch_master_placeholder(xml, "dt")
    return xml


def patch_ppt(path: Path, title: str) -> bytes:
    """Patch a presentation package in place; returns the new bytes."""
    order, parts = read_parts(path)
    parts["ppt/theme/theme1.xml"] = patch_theme(parts["ppt/theme/theme1.xml"].decode()).encode()
    master = "ppt/slideMasters/slideMaster1.xml"
    parts[master] = patch_ppt_master(parts[master].decode()).encode()
    for name in list(parts):
        if re.fullmatch(r"ppt/slideLayouts/slideLayout\d+\.xml", name):
            layout = parts[name].decode()
            layout = layout.replace('<a:buFont typeface="Arial"/>', '<a:buFont typeface="Segoe UI Variable Text"/>')
            layout = layout.replace('<a:latin typeface="Calibri"/>', '<a:latin typeface="+mn-lt"/>')
            parts[name] = layout.encode()
    parts["docProps/core.xml"] = patch_core(parts["docProps/core.xml"].decode(), title).encode()
    parts["docProps/app.xml"] = patch_app(parts["docProps/app.xml"].decode()).encode()
    data = zip_bytes(order, parts)
    emit(path, data)
    return data


# --- specialised decks (built FROM the patched deccan-deck.potx) --------------

_CT_PPT_TEMPLATE = b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
_CT_PPT_PRESENTATION = b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
_CT_WORD_TEMPLATE = b"application/vnd.openxmlformats-officedocument.wordprocessingml.template.main+xml"
_CT_WORD_DOCUMENT = b"application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"
_CT_XL_TEMPLATE = b"application/vnd.openxmlformats-officedocument.spreadsheetml.template.main+xml"
_CT_XL_SHEET = b"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"


def swap_content_type(package: bytes, old: bytes, new: bytes) -> bytes:
    order, parts = read_parts(package)
    ct = parts["[Content_Types].xml"]
    if old not in ct:
        raise SystemExit("expected content type not found during swap")
    parts["[Content_Types].xml"] = ct.replace(old, new)
    return zip_bytes(order, parts)


def build_specialized_deck(deck_bytes: bytes, spec: dict) -> bytes:
    """Rebuild the deck's sample slides for a specialised template.

    Opens the patched .potx as a presentation via content-type swap, replaces
    the sample slides with the specialisation's slide list (reusing the deck's
    sample-slide geometry: 0.6in margins, 56x2px blue rule, mono eyebrows),
    then swaps back to a template package deterministically.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    BLUE = RGBColor(0x16, 0x49, 0x99)
    S500 = RGBColor(0x78, 0x71, 0x6C)
    S700 = RGBColor(0x44, 0x40, 0x3C)
    S800 = RGBColor(0x29, 0x25, 0x24)
    S900 = RGBColor(0x1C, 0x19, 0x17)
    S100 = RGBColor(0xF5, 0xF5, 0xF4)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    pres = Presentation(io.BytesIO(swap_content_type(deck_bytes, _CT_PPT_TEMPLATE, _CT_PPT_PRESENTATION)))
    blank = next(l for l in pres.slide_masters[0].slide_layouts if l.name == "Blank")

    sld_ids = pres.slides._sldIdLst
    for sld_id in list(sld_ids):
        pres.part.drop_rel(sld_id.rId)
        sld_ids.remove(sld_id)

    def box(slide, left, top, w, h):
        b = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        b.text_frame.word_wrap = True
        return b

    def run(para, text, font, size, bold, color):
        r = para.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        return r

    def rule(slide, left, top):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.6), Inches(0.04))
        sh.fill.solid()
        sh.fill.fore_color.rgb = BLUE
        sh.line.fill.background()

    def title_bar(slide, text):
        run(box(slide, 0.6, 0.5, 12.0, 0.8).text_frame.paragraphs[0],
            text, T.SANS_TEXT, 24, True, S900)

    def bullets(slide, left, top, w, items, size=18):
        tf = box(slide, left, top, w, 5.0).text_frame
        for i, item in enumerate(items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run(para, item, T.SANS_TEXT, size, False, S800)
            para.space_after = Pt(8)

    def eyebrow(slide, text, left=0.6, top=0.4):
        run(box(slide, left, top, 4.0, 0.4).text_frame.paragraphs[0],
            text, T.MONO, 11, True, S500)

    # Slide builders keyed by kind.
    def cover(slide, s):
        run(box(slide, 0.6, 2.8, 11.0, 0.4).text_frame.paragraphs[0],
            "Deccan Fine Chemicals", T.SANS_TEXT, 12, True, S900)
        rule(slide, 0.6, 3.5)
        run(box(slide, 0.6, 3.7, 11.0, 1.4).text_frame.paragraphs[0],
            s["title"], T.SANS_DISPLAY, 44, False, BLUE)
        run(box(slide, 0.6, 5.0, 11.0, 0.6).text_frame.paragraphs[0],
            s["subtitle"], T.SANS_TEXT, 18, False, S700)
        run(box(slide, 0.6, 6.4, 8.0, 0.4).text_frame.paragraphs[0],
            "MONTH YYYY · VERSION 1.0 · CONFIDENTIAL", T.MONO, 10, True, S500)
        eyebrow(slide, s["eyebrow"], top=0.4)

    def section(slide, s):
        eyebrow(slide, s["eyebrow"])
        run(box(slide, 0.6, 3.0, 12.0, 1.4).text_frame.paragraphs[0],
            s["title"], T.SANS_DISPLAY, 44, False, BLUE)
        rule(slide, 0.6, 4.6)

    def onecol(slide, s):
        title_bar(slide, s["title"])
        bullets(slide, 0.6, 1.4, 12.0, s["items"])

    def twocol(slide, s):
        title_bar(slide, s["title"])
        for i, (head, items) in enumerate(s["columns"]):
            left = 0.6 + i * 6.4
            eyebrow(slide, head, left=left, top=1.4)
            bullets(slide, left, 2.0, 5.8, items, size=16)

    def threecol(slide, s):
        title_bar(slide, s["title"])
        for i, (head, items) in enumerate(s["columns"]):
            left = 0.6 + i * 4.1
            eyebrow(slide, head, left=left, top=1.4)
            bullets(slide, left, 2.0, 3.8, items, size=14)

    def table(slide, s):
        title_bar(slide, s["title"])
        data = s["rows"]
        rows, cols = len(data), len(data[0])
        shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.6),
                                       Inches(12.0), Inches(min(0.4 * rows, 4.8)))
        tbl = shape.table
        tbl.first_row = False
        tbl.horz_banding = False
        for r_i, row in enumerate(data):
            for c_i, text in enumerate(row):
                cell = tbl.cell(r_i, c_i)
                para = cell.text_frame.paragraphs[0]
                para.alignment = PP_ALIGN.LEFT
                cell.fill.solid()
                if r_i == 0:
                    run(para, text.upper(), T.MONO, 10, True, S700)
                    cell.fill.fore_color.rgb = S100
                else:
                    run(para, text, T.SANS_TEXT, 12, False, S800)
                    cell.fill.fore_color.rgb = WHITE

    def end(slide, s):
        width = pres.slide_width / 914400
        b = box(slide, 0.0, 3.2, width, 0.5)
        b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run(b.text_frame.paragraphs[0], "Deccan Fine Chemicals", T.SANS_TEXT, 21, True, S900)
        rule(slide, width / 2 - 0.3, 3.9)
        b = box(slide, 0.0, 4.2, width, 0.4)
        b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run(b.text_frame.paragraphs[0], "deccanchemicals.com · Hyderabad, India",
            T.SANS_TEXT, 12, False, S700)

    builders = {"cover": cover, "section": section, "onecol": onecol,
                "twocol": twocol, "threecol": threecol, "table": table, "end": end}
    for slide_spec in spec["slides"]:
        slide = pres.slides.add_slide(blank)
        builders[slide_spec["kind"]](slide, slide_spec)

    out = io.BytesIO()
    pres.save(out)
    as_template = swap_content_type(out.getvalue(), _CT_PPT_PRESENTATION, _CT_PPT_TEMPLATE)
    # Re-zip deterministically (python-pptx output is not timestamp-stable).
    order, parts = read_parts(as_template)
    parts["docProps/core.xml"] = patch_core(parts["docProps/core.xml"].decode(), spec["title"]).encode()
    parts["docProps/app.xml"] = patch_app(parts["docProps/app.xml"].decode()).encode()
    return zip_bytes(order, parts)


PITCH_SPEC = {
    "title": "Deccan Customer Pitch",
    "slides": [
        {"kind": "cover", "eyebrow": "CUSTOMER PITCH", "title": "Customer pitch title",
         "subtitle": "One sentence on the customer problem this pitch addresses."},
        {"kind": "section", "eyebrow": "SECTION 01", "title": "Why Deccan"},
        {"kind": "onecol", "title": "Why Deccan", "items": [
            "Custom chemical manufacturing at validated commercial scale.",
            "Regulatory track record: AEO-T1, QHSE audit history available under NDA.",
            "Dedicated program management from technology transfer to steady state.",
        ]},
        {"kind": "threecol", "title": "Capabilities", "columns": [
            ("CHEMISTRY", ["Multi-step synthesis.", "Hazardous chemistry handled at scale."]),
            ("CAPACITY", ["Reactor volumes and materials of construction on request."]),
            ("COMPLIANCE", ["QHSE systems.", "Customer and regulatory audits supported."]),
        ]},
        {"kind": "table", "title": "Proof points", "rows": [
            ["Program", "Scale", "Outcome", "Reference"],
            ["Replace with program", "Replace", "Replace", "Available on request"],
            ["Replace with program", "Replace", "Replace", "Available on request"],
        ]},
        {"kind": "twocol", "title": "Next steps", "columns": [
            ("DECCAN", ["Technical questionnaire response.", "Site visit agenda."]),
            ("CUSTOMER", ["Process package under CDA.", "Target timeline confirmation."]),
        ]},
        {"kind": "end"},
    ],
}

REVIEW_SPEC = {
    "title": "Deccan Internal Review",
    "slides": [
        {"kind": "cover", "eyebrow": "INTERNAL REVIEW", "title": "Internal review title",
         "subtitle": "Reporting period and scope in one sentence."},
        {"kind": "onecol", "title": "Agenda", "items": [
            "Performance against plan.",
            "Key risks and mitigations.",
            "Decisions needed from this review.",
            "Actions and owners.",
        ]},
        {"kind": "table", "title": "KPIs", "rows": [
            ["Metric", "Target", "Actual", "Trend"],
            ["Replace with metric", "Replace", "Replace", "Replace"],
            ["Replace with metric", "Replace", "Replace", "Replace"],
        ]},
        {"kind": "twocol", "title": "Decisions needed", "columns": [
            ("DECISION", ["Replace with the decision to be taken."]),
            ("CONTEXT", ["Replace with the constraint or trade-off."]),
        ]},
        {"kind": "table", "title": "Actions and owners", "rows": [
            ["Action", "Owner", "Due", "Status"],
            ["Replace with action", "Replace", "Replace", "Open"],
            ["Replace with action", "Replace", "Replace", "Open"],
        ]},
        {"kind": "end"},
    ],
}


# --- Google Workspace derivation ----------------------------------------------


def derive_gworkspace() -> None:
    gw = TEMPLATES / "gworkspace"
    emit(
        gw / "deccan-document-for-drive.docx",
        swap_content_type((TEMPLATES / "word" / "deccan-document.dotx").read_bytes(),
                          _CT_WORD_TEMPLATE, _CT_WORD_DOCUMENT),
    )
    emit(
        gw / "deccan-workbook-for-drive.xlsx",
        swap_content_type((TEMPLATES / "excel" / "deccan-workbook.xltx").read_bytes(),
                          _CT_XL_TEMPLATE, _CT_XL_SHEET),
    )
    emit(
        gw / "deccan-deck-for-drive.pptx",
        swap_content_type((TEMPLATES / "powerpoint" / "deccan-deck.potx").read_bytes(),
                          _CT_PPT_TEMPLATE, _CT_PPT_PRESENTATION),
    )


# --- main ----------------------------------------------------------------------


def main() -> None:
    global CHECK_MODE
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--check", action="store_true", help="verify outputs, do not write")
    CHECK_MODE = parser.parse_args().check

    print("Word templates:")
    for name, title in WORD_TITLES.items():
        patch_word(TEMPLATES / "word" / name, title)

    print("Excel templates:")
    for name, title in EXCEL_TITLES.items():
        patch_excel(TEMPLATES / "excel" / name, title)

    print("PowerPoint templates:")
    deck_path = TEMPLATES / "powerpoint" / "deccan-deck.potx"
    # Use the returned bytes, not a re-read: in --check mode the on-disk deck
    # may be unpatched, and the specialised decks derive from the patched one.
    deck_bytes = patch_ppt(deck_path, "Deccan Deck")
    emit(TEMPLATES / "powerpoint" / "deccan-customer-pitch.potx",
         build_specialized_deck(deck_bytes, PITCH_SPEC))
    emit(TEMPLATES / "powerpoint" / "deccan-internal-review.potx",
         build_specialized_deck(deck_bytes, REVIEW_SPEC))

    print("Google Workspace derivations:")
    derive_gworkspace()

    if CHECK_MODE and _changed:
        print(f"\n{len(_changed)} file(s) out of date. Run: python tools/templates-build/build_templates.py")
        sys.exit(1)
    print("done.")


if __name__ == "__main__":
    main()
